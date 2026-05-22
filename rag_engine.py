import os
os.environ['OPENBLAS_NUM_THREADS'] = '1'
os.environ['OMP_NUM_THREADS'] = '1'
os.environ['MKL_NUM_THREADS'] = '1'

import shutil
import re
import uuid
import json
from typing import List, Dict, Optional
from config import (DOCUMENT_FOLDER, VECTOR_DB_PATH, CHUNK_SIZE, CHUNK_OVERLAP,
                    TOP_K, LLM_ENABLED, LLM_TYPE, LLM_MODEL, LLM_API_KEY,
                    LLM_BASE_URL, LLM_LOCAL_URL, LLM_LOCAL_MODEL,
                    LLM_CONFIG_FILE, LLM_PRESETS)


class RAGEngine:
    def __init__(self):
        self.initialized = False
        self.added_sources = {}
        self.doc_cache = []
        self.cache_version = -1
        self.llm_enabled = LLM_ENABLED
        self.llm_type = LLM_TYPE
        self.llm_model = LLM_MODEL
        self.llm_api_key = LLM_API_KEY
        self.llm_base_url = LLM_BASE_URL
        self.llm_local_url = LLM_LOCAL_URL
        self.llm_local_model = LLM_LOCAL_MODEL
        self._load_llm_config()

    def initialize(self, device: str = "cpu"):
        import chromadb

        self.chroma_client = chromadb.PersistentClient(path=VECTOR_DB_PATH)

        try:
            self.collection = self.chroma_client.get_collection("knowledge_base")
        except Exception:
            self.collection = self.chroma_client.create_collection("knowledge_base")

        self.initialized = True

        count = self.collection.count()
        if count > 0:
            self._load_sources_from_db()
        self._refresh_cache()
        print(f"引擎就绪，知识库: {count} 条片段")

    def _load_sources_from_db(self):
        try:
            count = self.collection.count()
            data = self.collection.get(include=["metadatas"], limit=count)
            metas = data.get("metadatas", [])
            self.added_sources = {}
            for meta in metas:
                src = meta.get("source", "")
                if src:
                    if src not in self.added_sources:
                        self.added_sources[src] = {"source": src, "chunk_count": 0}
                    self.added_sources[src]["chunk_count"] += 1
        except Exception:
            self.added_sources = {}

    def _invalidate_cache(self):
        self.cache_version = -1

    def _refresh_cache(self):
        try:
            count = self.collection.count()
        except Exception:
            return
        if self.cache_version == count and self.doc_cache:
            return
        try:
            data = self.collection.get(include=["documents", "metadatas"], limit=count)
            docs = data.get("documents", [])
            metas = data.get("metadatas", [])
            self.doc_cache = [(docs[i], metas[i] if i < len(metas) else {}) for i in range(len(docs))]
            self.cache_version = len(self.doc_cache)
        except Exception:
            pass

    def _read_file(self, file_path: str) -> Optional[str]:
        try:
            if file_path.endswith(".txt"):
                with open(file_path, "r", encoding="utf-8") as f:
                    return f.read()
            elif file_path.endswith(".pdf"):
                import fitz
                doc = fitz.open(file_path)
                return "".join(page.get_text() for page in doc)
            elif file_path.endswith(".docx"):
                import docx
                doc = docx.Document(file_path)
                return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            elif file_path.endswith(".pptx"):
                import pptx
                prs = pptx.Presentation(file_path)
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                t = para.text.strip()
                                if t:
                                    texts.append(t)
                        if shape.has_table:
                            table = shape.table
                            for row in table.rows:
                                row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                                if row_texts:
                                    texts.append(" | ".join(row_texts))
                return "\n".join(texts)
            elif file_path.endswith(".xlsx"):
                import openpyxl
                wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                texts = []
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    texts.append(f"[Sheet: {sheet_name}]")
                    row_count = 0
                    for row in ws.iter_rows(values_only=True):
                        row_texts = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                        if row_texts:
                            texts.append(" | ".join(row_texts))
                            row_count += 1
                            if row_count > 500:
                                texts.append("... (truncated)")
                                break
                    texts.append("")
                wb.close()
                return "\n".join(texts)
            return None
        except Exception as e:
            print(f"  读取失败 {os.path.basename(file_path)}: {e}")
            return None

    def _split_text(self, text: str, source: str) -> List[Dict]:
        chunks = []
        text = text.replace("\r\n", "\n")
        text_len = len(text)
        pos = 0
        chunk_id = 0
        max_chunks = 5000
        prev_pos = -1

        while pos < text_len and chunk_id < max_chunks:
            if pos == prev_pos:
                pos += CHUNK_SIZE
            prev_pos = pos

            cut = min(pos + CHUNK_SIZE, text_len)
            for sep in ["\n\n", "\n", "。", "！", "？"]:
                idx = text.rfind(sep, pos, cut)
                if idx > pos + CHUNK_SIZE // 2:
                    cut = idx + len(sep)
                    break

            chunk_text = text[pos:cut].strip()
            if chunk_text and len(chunk_text) > 10:
                chunks.append({"content": chunk_text, "source": source, "chunk_id": chunk_id})
                chunk_id += 1

            next_pos = cut - CHUNK_OVERLAP
            if next_pos <= pos:
                next_pos = cut
            pos = next_pos

        return chunks

    def load_documents(self, folder_path: str = DOCUMENT_FOLDER) -> List[Dict]:
        documents = []
        if not os.path.exists(folder_path):
            os.makedirs(folder_path)
            return documents
        for filename in os.listdir(folder_path):
            file_path = os.path.join(folder_path, filename)
            content = self._read_file(file_path)
            if content:
                chunks = self._split_text(content, file_path)
                documents.extend(chunks)
        return documents

    def load_single_document(self, file_path: str) -> List[Dict]:
        content = self._read_file(file_path)
        if content is None:
            return []
        return self._split_text(content, file_path)

    def build_knowledge_base(self, folder_path: str = DOCUMENT_FOLDER) -> Dict:
        if not self.initialized:
            return {"status": "error", "message": "引擎未初始化"}

        docs = self.load_documents(folder_path)
        if not docs:
            return {"status": "error", "message": "文件夹中没有支持的文档 (PDF/TXT/DOCX/PPTX/XLSX)"}

        try:
            self.chroma_client.delete_collection("knowledge_base")
        except Exception:
            pass
        self.collection = self.chroma_client.create_collection("knowledge_base")

        self.added_sources = {}
        for d in docs:
            src = d["source"]
            if src not in self.added_sources:
                self.added_sources[src] = {"source": src, "chunk_count": 0}
            self.added_sources[src]["chunk_count"] += 1

        ids = [str(uuid.uuid4()) for _ in docs]
        metadatas = [{"source": d["source"], "chunk_id": d["chunk_id"]} for d in docs]
        documents = [d["content"] for d in docs]

        print(f"正在写入 {len(documents)} 个文档片段...")
        self.collection.add(ids=ids, metadatas=metadatas, documents=documents)
        self._invalidate_cache()

        return {"status": "success", "document_count": len(set(d["source"] for d in docs)), "chunk_count": len(docs)}

    def add_documents(self, file_paths: List[str]) -> Dict:
        if not self.initialized:
            return {"status": "error", "message": "引擎未初始化"}

        all_docs = []
        for fp in file_paths:
            all_docs.extend(self.load_single_document(fp))

        if not all_docs:
            return {"status": "error", "message": "没有有效文档"}

        for d in all_docs:
            src = d["source"]
            if src not in self.added_sources:
                self.added_sources[src] = {"source": src, "chunk_count": 0}
            self.added_sources[src]["chunk_count"] += 1

        ids = [str(uuid.uuid4()) for _ in all_docs]
        metadatas = [{"source": d["source"], "chunk_id": d["chunk_id"]} for d in all_docs]
        documents = [d["content"] for d in all_docs]

        print(f"正在追加 {len(documents)} 个文档片段...")
        try:
            self.collection.add(ids=ids, metadatas=metadatas, documents=documents)
        except Exception as e:
            return {"status": "error", "message": f"写入失败: {str(e)}"}

        self._invalidate_cache()
        return {"status": "success", "chunk_count": len(all_docs)}

    def _extract_keywords(self, text: str) -> List[str]:
        words = re.findall(r'[\u4e00-\u9fff]{2,}|[a-zA-Z]{2,}', text)
        stop_words = {'的', '了', '在', '是', '我', '有', '和', '就', '不', '人', '都', '一', '一个', '上',
                      '也', '很', '到', '说', '要', '去', '你', '会', '着', '没有', '看', '好', '自己',
                      '这', '他', '她', '它', '们', '那', '什么', '怎么', '如何', '可以', '这个', '那个',
                      'the', 'a', 'an', 'is', 'are', 'was', 'were', 'be', 'been', 'do', 'does', 'did',
                      'this', 'that', 'these', 'those', 'to', 'of', 'in', 'for', 'on', 'with', 'and'}
        keywords = [w for w in words if w.lower() not in stop_words]

        expanded = list(keywords)
        for kw in keywords:
            if len(kw) > 2:
                for i in range(len(kw) - 1):
                    bigram = kw[i:i + 2]
                    if bigram.lower() not in stop_words and len(bigram) >= 2:
                        expanded.append(bigram)

        result = list(dict.fromkeys(expanded))
        if not result:
            chars = re.findall(r'[\u4e00-\u9fff]|[a-zA-Z]+', text)
            result = [c for c in chars if c.lower() not in stop_words]
        return result[:30]

    def _score_chunk(self, keywords: List[str], chunk: str) -> float:
        score = 0.0
        chunk_lower = chunk.lower()
        for kw in keywords:
            count = chunk_lower.count(kw.lower())
            if count > 0:
                score += count * (1.0 + len(kw) * 0.5)
        return score

    def query(self, question: str, top_k: int = TOP_K) -> Dict:
        if not self.initialized:
            return {"status": "error", "message": "引擎未初始化"}

        self._refresh_cache()

        if not self.doc_cache:
            return {"status": "error", "message": "知识库为空，请先构建"}

        keywords = self._extract_keywords(question)

        scored = []
        for i, (doc, meta) in enumerate(self.doc_cache):
            score = self._score_chunk(keywords, doc)
            if score > 0:
                scored.append((score, i))

        if not scored:
            fallback_chars = re.findall(r'[\u4e00-\u9fff]|[a-zA-Z]+', question)
            stop_words = {'的', '了', '在', '是', '我', '有', '和', '就', '不', '人', '都', '一', '一个', '上',
                          '也', '很', '到', '说', '要', '去', '你', '会', '着', '没有', '看', '好', '自己',
                          '这', '他', '她', '它', '们', '那', '什么', '怎么', '如何', '可以', '这个', '那个',
                          'the', 'a', 'an', 'is', 'are', 'was', 'were', 'be', 'been', 'do', 'does', 'did',
                          'this', 'that', 'these', 'those', 'to', 'of', 'in', 'for', 'on', 'with', 'and'}
            fallback_chars = [c for c in fallback_chars if c.lower() not in stop_words]
            fallback_chars = list(dict.fromkeys(fallback_chars))[:10]
            if fallback_chars:
                for i, (doc, meta) in enumerate(self.doc_cache):
                    score = self._score_chunk(fallback_chars, doc)
                    if score > 0:
                        scored.append((score * 0.3, i))

        scored.sort(key=lambda x: x[0], reverse=True)
        top_results = scored[:top_k]

        best_score = top_results[0][0] if top_results else 0
        if not top_results or best_score < 5.0:
            print(f"[Query] 本地搜索结果不足 (最高分={best_score:.1f}), llm_enabled={self.llm_enabled}, type={self.llm_type}, has_key={bool(self.llm_api_key)}")
            if self.llm_enabled and self.llm_type == "online" and self.llm_api_key:
                answer = self._generate_answer_from_web(question)
                return {
                    "status": "success",
                    "answer": answer,
                    "sources": [],
                    "from_web": True
                }
            return {"status": "error", "message": "未找到相关内容，请尝试更换关键词"}

        sources = []
        for score, idx in top_results:
            doc, meta = self.doc_cache[idx]
            sources.append({
                "content": doc,
                "source": meta.get("source", "unknown"),
                "score": score
            })

        if self.llm_enabled and self.llm_type == "online" and self.llm_api_key and best_score < 20.0:
            print(f"[Query] 本地匹配较弱 (最高分={best_score:.1f}), 增强联网搜索")
            web_info = self._search_web(question)
            if web_info:
                answer = self._generate_answer_llm_with_web(question, sources, web_info)
                return {
                    "status": "success",
                    "answer": answer,
                    "sources": sources,
                    "from_web": True
                }

        answer = self._generate_answer(question, sources, keywords) if not self.llm_enabled else self._generate_answer_llm(question, sources, keywords)
        return {"status": "success", "answer": answer, "sources": sources}

    def _generate_answer(self, question: str, sources: List[Dict], keywords: List[str]) -> str:
        answer_parts = []

        for i, src in enumerate(sources):
            content = src["content"]
            sentences = re.split(r'[。！？\n]', content)
            sentences = [s.strip() for s in sentences if len(s.strip()) > 5]

            if sentences:
                if i == 0 and len(sentences) >= 2:
                    best_start = 0
                    best_hits = 0
                    for si in range(len(sentences) - 1):
                        hits = sum(1 for kw in keywords if kw in sentences[si] + sentences[si + 1])
                        if hits > best_hits:
                            best_hits = hits
                            best_start = si
                    s1 = sentences[best_start][:300]
                    if best_start + 1 < len(sentences):
                        s2 = sentences[best_start + 1][:300]
                        answer_parts.append(f"{s1}。{s2}")
                    else:
                        answer_parts.append(s1)
                else:
                    best = max(sentences, key=lambda s: sum(1 for kw in keywords if kw in s))
                    if best not in answer_parts:
                        answer_parts.append(best[:300])

        answer = "。\n".join(answer_parts[:3])

        source_files = list(dict.fromkeys(os.path.basename(s["source"]) for s in sources))
        answer += f"\n\n参考文件: {', '.join(source_files[:3])}"

        return answer

    def _load_llm_config(self):
        try:
            if os.path.exists(LLM_CONFIG_FILE):
                with open(LLM_CONFIG_FILE, "r", encoding="utf-8") as f:
                    saved = json.load(f)
                if saved.get("enabled") is not None:
                    self.llm_enabled = saved["enabled"]
                if saved.get("type"):
                    self.llm_type = saved["type"]
                if saved.get("model"):
                    self.llm_model = saved["model"]
                if saved.get("api_key"):
                    self.llm_api_key = saved["api_key"]
                if saved.get("base_url"):
                    self.llm_base_url = saved["base_url"]
                if saved.get("local_url"):
                    self.llm_local_url = saved["local_url"]
                if saved.get("local_model"):
                    self.llm_local_model = saved["local_model"]
                print(f"[LLM] 已加载配置: enabled={self.llm_enabled}, type={self.llm_type}, model={self.llm_model}, base_url={self.llm_base_url}")
        except Exception as e:
            print(f"[LLM] 加载配置失败: {e}")

    def _save_llm_config(self):
        try:
            saved = {
                "enabled": self.llm_enabled,
                "type": self.llm_type,
                "model": self.llm_model,
                "api_key": self.llm_api_key,
                "base_url": self.llm_base_url,
                "local_url": self.llm_local_url,
                "local_model": self.llm_local_model,
            }
            with open(LLM_CONFIG_FILE, "w", encoding="utf-8") as f:
                json.dump(saved, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"[LLM] 保存配置失败: {e}")

    def update_llm_config(self, config: Dict):
        if "enabled" in config:
            self.llm_enabled = config["enabled"]
        if "type" in config:
            self.llm_type = config["type"]
        if "model" in config:
            self.llm_model = config["model"]
        if "api_key" in config and config["api_key"]:
            self.llm_api_key = config["api_key"]
        if "base_url" in config:
            self.llm_base_url = config["base_url"]
        if "local_url" in config:
            self.llm_local_url = config["local_url"]
        if "local_model" in config:
            self.llm_local_model = config["local_model"]
        self._save_llm_config()
        print(f"[LLM] 配置已保存: enabled={self.llm_enabled}, type={self.llm_type}, model={self.llm_model}")

    def get_llm_config(self) -> Dict:
        return {
            "enabled": self.llm_enabled,
            "type": self.llm_type,
            "model": self.llm_model,
            "api_key": "***" if self.llm_api_key else "",
            "has_api_key": bool(self.llm_api_key),
            "base_url": self.llm_base_url,
            "local_url": self.llm_local_url,
            "local_model": self.llm_local_model
        }

    def get_llm_presets(self) -> List[Dict]:
        return list(LLM_PRESETS)

    def _generate_answer_llm(self, question: str, sources: List[Dict],
                             keywords: List[str]) -> str:
        context = ""
        for i, src in enumerate(sources[:5]):
            context += f"[来源{i + 1}: {os.path.basename(src['source'])}]\n{src['content'][:600]}\n\n"

        prompt = (
            "你是一个知识库问答助手。请按以下优先级回答用户问题：\n"
            "1）如果参考资料与问题相关，优先基于资料回答；\n"
            "2）如果参考资料与问题无关，请根据你自身的知识直接回答。\n\n"
            f"=== 参考资料 ===\n{context}\n"
            f"=== 用户问题 ===\n{question}\n\n"
            "请用中文回答："
        )

        if self.llm_type == "local":
            return self._call_local_llm(prompt, sources)
        else:
            return self._call_online_llm(prompt, sources)

    def _call_online_llm(self, prompt: str, sources: List[Dict]) -> str:
        try:
            import urllib.request
            import urllib.error

            payload = json.dumps({
                "model": self.llm_model,
                "messages": [
                    {"role": "system", "content": "你是一个智能问答助手，基于提供的参考资料或实时信息如实回答。"},
                    {"role": "user", "content": prompt}
                ],
                "temperature": 0.3,
                "max_tokens": 2000
            }).encode("utf-8")

            url = f"{self.llm_base_url}/chat/completions"
            print(f"[LLM] 调用在线模型: {url}")
            print(f"[LLM] 模型: {self.llm_model}, Key: {self.llm_api_key[:8]}...")
            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {self.llm_api_key}"
            }
            req = urllib.request.Request(url, data=payload, headers=headers, method="POST")

            with urllib.request.urlopen(req, timeout=60) as resp:
                result = json.loads(resp.read().decode("utf-8"))
                answer = result["choices"][0]["message"]["content"]
                print(f"[LLM] 在线模型返回成功，长度: {len(answer)}")
        except urllib.error.HTTPError as e:
            body = e.read().decode("utf-8", errors="replace")
            print(f"[LLM] HTTP错误 {e.code}: {body}")
            return f"[大模型调用失败: HTTP {e.code}]\n\n{body[:300]}\n\n已切换为原文摘抄模式。\n\n" + self._generate_answer_legacy(sources)
        except Exception as e:
            print(f"[LLM] 异常: {e}")
            return f"[大模型调用失败: {e}]\n\n已切换为原文摘抄模式。\n\n" + self._generate_answer_legacy(sources)

        source_files = list(dict.fromkeys(os.path.basename(s["source"]) for s in sources))
        answer += f"\n\n📄 参考文件: {', '.join(source_files[:5])}"
        return answer

    def _call_local_llm(self, prompt: str, sources: List[Dict]) -> str:
        try:
            import urllib.request
            import urllib.error

            payload = json.dumps({
                "model": self.llm_local_model,
                "prompt": prompt,
                "stream": False,
                "options": {"temperature": 0.3}
            }).encode("utf-8")

            req = urllib.request.Request(
                self.llm_local_url,
                data=payload,
                headers={"Content-Type": "application/json"},
                method="POST"
            )

            with urllib.request.urlopen(req, timeout=120) as resp:
                result = json.loads(resp.read().decode("utf-8"))
                answer = result.get("response", "")
        except Exception as e:
            return f"[本地模型调用失败: {e}]\n\n请确认 Ollama 已启动且模型已下载。\n\n" + self._generate_answer_legacy(sources)

        source_files = list(dict.fromkeys(os.path.basename(s["source"]) for s in sources))
        answer += f"\n\n📄 参考文件: {', '.join(source_files[:5])}"
        return answer

    def _generate_answer_llm_with_web(self, question: str, sources: List[Dict],
                                       web_info: str) -> str:
        context = ""
        for i, src in enumerate(sources[:3]):
            context += f"[本地{i + 1}: {os.path.basename(src['source'])}]\n{src['content'][:400]}\n\n"

        prompt = (
            "你是一个知识库问答助手。以下是本地知识库和实时网络搜索的结果，请综合回答用户问题。\n"
            "优先使用实时网络信息回答，本地资料作为补充参考。\n\n"
            f"=== 本地资料 ===\n{context}\n"
            f"=== 实时网络搜索结果 ===\n{web_info}\n\n"
            f"=== 用户问题 ===\n{question}\n\n"
            "请用中文回答："
        )

        return self._call_online_llm(prompt, sources)

    def _generate_answer_legacy(self, sources: List[Dict]) -> str:
        answer_parts = []
        for src in sources[:3]:
            sentences = re.split(r'[。！？\n]', src["content"])
            sentences = [s.strip() for s in sentences if len(s.strip()) > 5]
            if sentences:
                answer_parts.append(sentences[0][:300])

        answer = "。\n".join(answer_parts)
        source_files = list(dict.fromkeys(os.path.basename(s["source"]) for s in sources))
        answer += f"\n\n📄 参考文件: {', '.join(source_files[:5])}"
        return answer

    def _search_web(self, query: str) -> str:
        try:
            import urllib.request
            import urllib.parse
            import html as html_module

            search_url = "https://lite.duckduckgo.com/lite/?" + urllib.parse.urlencode({"q": query})
            headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"}
            req = urllib.request.Request(search_url, headers=headers)

            with urllib.request.urlopen(req, timeout=10) as resp:
                raw = resp.read().decode("utf-8", errors="replace")

            results = []
            for match in re.finditer(r'<a[^>]*class="result-link"[^>]*href="([^"]*)"[^>]*>([^<]*)</a>', raw):
                url = match.group(1)
                title = match.group(2).strip()
                results.append({"title": title, "url": url, "snippet": ""})

            snippets = re.findall(r'<td class="result-snippet"[^>]*>(.*?)</td>', raw, re.DOTALL)
            for i, s in enumerate(snippets):
                if i < len(results):
                    clean = re.sub(r'<[^>]+>', '', s).strip()
                    results[i]["snippet"] = html_module.unescape(clean)

            print(f"[WebSearch] DuckDuckGo 搜索 '{query}' 得到 {len(results)} 条结果")
            if results:
                lines = []
                for r in results[:5]:
                    lines.append(f"- [{r['title']}]({r['url']})\n  {r['snippet']}")
                return "\n\n".join(lines)
            return ""
        except Exception as e:
            print(f"[WebSearch] 搜索失败: {e}")
            return ""

    def _generate_answer_from_web(self, question: str) -> str:
        web_results = self._search_web(question)

        if web_results:
            prompt = (
                "你是一个智能问答助手。以下是实时搜索到的网页信息，请根据这些信息回答用户问题。\n"
                "要求：基于搜索结果作答，引用相关信息。如果搜索结果不足以回答问题，请如实说明。\n\n"
                f"=== 搜索结果 ===\n{web_results}\n\n"
                f"=== 用户问题 ===\n{question}\n\n"
                "请用中文回答："
            )
        else:
            prompt = (
                "你是一个智能问答助手。用户的问题在本地知识库中未找到匹配内容，"
                "且网络搜索暂时不可用。请根据你自身的知识回答以下问题。\n"
                "要求：如实作答，如果确实不知道请说明，不要编造信息。\n\n"
                f"=== 用户问题 ===\n{question}\n\n"
                "请用中文回答："
            )

        try:
            import urllib.request
            import urllib.error

            payload = json.dumps({
                "model": self.llm_model,
                "messages": [
                    {"role": "system", "content": "你是一个智能问答助手，基于实时信息或自身知识如实回答用户问题。"},
                    {"role": "user", "content": prompt}
                ],
                "temperature": 0.5,
                "max_tokens": 2000
            }).encode("utf-8")

            url = f"{self.llm_base_url}/chat/completions"
            print(f"[Web] 联网回答调用: {url}")
            print(f"[Web] 模型: {self.llm_model}")
            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {self.llm_api_key}"
            }
            req = urllib.request.Request(url, data=payload, headers=headers, method="POST")

            with urllib.request.urlopen(req, timeout=60) as resp:
                result = json.loads(resp.read().decode("utf-8"))
                answer = result["choices"][0]["message"]["content"]
                print(f"[Web] 联网回答成功，长度: {len(answer)}")
        except urllib.error.HTTPError as e:
            body = e.read().decode("utf-8", errors="replace")
            print(f"[Web] HTTP错误 {e.code}: {body}")
            return f"[联网检索失败: HTTP {e.code}]\n\n{body[:300]}"
        except Exception as e:
            print(f"[Web] 异常: {e}")
            return f"[联网检索失败: {e}]\n\n请检查网络连接和大模型 API 配置。"

        if web_results:
            answer += "\n\n🌐 回答来源: 实时网络搜索 + 大模型总结"
        else:
            answer += "\n\n🌐 回答来源: 大模型自身知识（网络搜索不可用）"
        return answer

    def get_db_stats(self) -> Dict:
        if not self.initialized:
            return {"document_count": 0, "chunk_count": 0}
        try:
            count = self.collection.count()
            return {"document_count": count, "chunk_count": count}
        except Exception:
            return {"document_count": 0, "chunk_count": 0}

    def get_added_files(self) -> List[Dict]:
        result = []
        for src, info in self.added_sources.items():
            entry = dict(info)
            if os.path.isfile(src):
                entry["bytes"] = os.path.getsize(src)
                entry["exists"] = True
            else:
                entry["bytes"] = 0
                entry["exists"] = False
            result.append(entry)
        return result

    def remove_document(self, source_path: str) -> Dict:
        if not self.initialized:
            return {"status": "error", "message": "引擎未初始化"}
        if source_path not in self.added_sources:
            return {"status": "error", "message": "该文件未在知识库中"}

        try:
            data = self.collection.get(include=["metadatas"])
            all_ids = data.get("ids", [])
            all_metas = data.get("metadatas", [])

            ids_to_delete = []
            for i, meta in enumerate(all_metas):
                if meta.get("source", "") == source_path:
                    ids_to_delete.append(all_ids[i])

            if ids_to_delete:
                self.collection.delete(ids=ids_to_delete)

            count = self.added_sources.pop(source_path, {}).get("chunk_count", 0)
            self._invalidate_cache()
            return {"status": "success", "message": f"已移除 {count} 个片段", "chunk_count": count}
        except Exception as e:
            return {"status": "error", "message": f"移除失败: {str(e)}"}

    def clear_knowledge_base(self) -> Dict:
        try:
            self.chroma_client.delete_collection("knowledge_base")
        except Exception:
            pass
        if os.path.exists(VECTOR_DB_PATH):
            try:
                shutil.rmtree(VECTOR_DB_PATH)
            except Exception:
                pass
        self.added_sources = {}
        self.doc_cache = []
        self.cache_version = -1
        self.chroma_client = None
        return {"status": "success", "message": "知识库已清空"}